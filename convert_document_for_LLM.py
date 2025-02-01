import os
import mimetypes
import subprocess
from PIL import Image
import pytesseract
import aspose.words as aw
from pptx import Presentation
from ebooklib import epub
from bs4 import BeautifulSoup
import pandas as pd

def detect_file_type(file_path):
    mime_type, _ = mimetypes.guess_type(file_path)
    return mime_type

def convert_pptx_to_text(file_path):
    text = []
    prs = Presentation(file_path)
    for i, slide in enumerate(prs.slides, 1):
        text.append(f"[SLIDE {i}]")
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n".join(text)

def convert_docx_to_text(file_path):
    doc = aw.Document(file_path)
    return doc.get_text()

def convert_xlsx_to_text(file_path):
    df = pd.read_excel(file_path, sheet_name=None)
    text = []
    for sheet_name, sheet_data in df.items():
        text.append(f"[SHEET: {sheet_name}]")
        text.append(sheet_data.to_string(index=False))
    return "\n\n".join(text)

def convert_epub_to_text(file_path):
    book = epub.read_epub(file_path)
    text = []
    for item in book.get_items():
        if item.get_type() == epub.ITEM_DOCUMENT:
            soup = BeautifulSoup(item.get_content(), 'html.parser')
            text.append(soup.get_text())
    return "\n\n".join(text)

def convert_image_to_text(file_path):
    return pytesseract.image_to_string(Image.open(file_path))

def convert_to_text(input_file, output_file):
    file_type = detect_file_type(input_file)
    
    if file_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
        text = convert_pptx_to_text(input_file)
    elif file_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
        text = convert_docx_to_text(input_file)
    elif file_type == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet":
        text = convert_xlsx_to_text(input_file)
    elif file_type == "application/epub+zip":
        text = convert_epub_to_text(input_file)
    elif file_type.startswith("image/"):
        text = convert_image_to_text(input_file)
    else:
        raise ValueError(f"Unsupported file type: {file_type}")

    with open(output_file, 'w', encoding='utf-8') as f:
        f.write(text)

    print(f"Conversion complete. Text saved to {output_file}")

# Example usage
input_file = "/path/to/your/input/file.pptx"  # Replace with your input file path
output_file = "/path/to/your/output/file.txt"  # Replace with your desired output file path

convert_to_text(input_file, output_file)
